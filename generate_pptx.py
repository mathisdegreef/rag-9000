"""Generate a PowerPoint presentation for rag-9000."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Colour palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # very dark navy
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # soft blue
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # soft green
ACCENT3    = RGBColor(0xF3, 0x8B, 0xA8)   # soft rose
ACCENT4    = RGBColor(0xF9, 0xE2, 0xAF)   # soft yellow
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCB, 0xD7)
MID_GREY   = RGBColor(0x89, 0x8A, 0x9A)
BOX_BG     = RGBColor(0x31, 0x31, 0x4A)   # slightly lighter box
# ───────────────────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ─── Low-level helpers ─────────────────────────────────────────────────────────

def fill_slide_bg(slide, color=DARK_BG):
    """Solid background rectangle (covers whole slide)."""
    left = top = Inches(0)
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.zorder = 0
    return shape


def add_rect(slide, l, t, w, h, fill=BOX_BG, line_color=None, radius=False):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_accent_bar(slide, color=ACCENT, height=Inches(0.07)):
    """Thin horizontal accent line near top."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0.75), SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def slide_title(slide, title, subtitle=None,
                title_color=ACCENT, sub_color=LIGHT_GREY):
    add_accent_bar(slide)
    add_text(slide, title,
             Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
             font_size=26, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
                 font_size=14, color=sub_color, italic=True)


def bullet_block(slide, items, l, t, w, h,
                 font_size=13, color=WHITE, bullet="▸ ", line_gap=0.38):
    """Render a list of strings as bullet points (one textbox per item)."""
    y = t
    step = Inches(line_gap)
    for item in items:
        add_text(slide, bullet + item, l, y, w, step,
                 font_size=font_size, color=color)
        y += step


def two_column_bullets(slide, left_items, right_items,
                       l_label=None, r_label=None,
                       top=Inches(1.5), font_size=12.5):
    col_w = Inches(5.9)
    gap   = Inches(0.5)
    l_x   = Inches(0.5)
    r_x   = l_x + col_w + gap

    if l_label:
        add_rect(slide, l_x, top - Inches(0.38), col_w, Inches(0.34),
                 fill=ACCENT, line_color=None)
        add_text(slide, l_label, l_x + Inches(0.1), top - Inches(0.38),
                 col_w - Inches(0.2), Inches(0.34),
                 font_size=12, bold=True, color=DARK_BG)
    if r_label:
        add_rect(slide, r_x, top - Inches(0.38), col_w, Inches(0.34),
                 fill=ACCENT3, line_color=None)
        add_text(slide, r_label, r_x + Inches(0.1), top - Inches(0.38),
                 col_w - Inches(0.2), Inches(0.34),
                 font_size=12, bold=True, color=DARK_BG)

    for i, item in enumerate(left_items):
        y = top + Inches(i * 0.43)
        add_rect(slide, l_x, y, col_w, Inches(0.38), fill=BOX_BG)
        add_text(slide, "▸  " + item, l_x + Inches(0.12), y + Inches(0.04),
                 col_w - Inches(0.2), Inches(0.34),
                 font_size=font_size, color=WHITE)

    for i, item in enumerate(right_items):
        y = top + Inches(i * 0.43)
        add_rect(slide, r_x, y, col_w, Inches(0.38), fill=BOX_BG)
        add_text(slide, "▸  " + item, r_x + Inches(0.12), y + Inches(0.04),
                 col_w - Inches(0.2), Inches(0.34),
                 font_size=font_size, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_slide_bg(slide)

    # big gradient-feel shape behind title
    banner = slide.shapes.add_shape(1, Inches(0), Inches(2.6), SLIDE_W, Inches(2.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = BOX_BG
    banner.line.fill.background()

    # decorative side stripe
    stripe = slide.shapes.add_shape(1, Inches(0), Inches(2.6), Inches(0.18), Inches(2.8))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    add_text(slide, "RAG-9000", Inches(0.4), Inches(2.7), Inches(12), Inches(1.0),
             font_size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "Retrieval Evaluation Framework",
             Inches(0.4), Inches(3.65), Inches(12), Inches(0.6),
             font_size=22, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Architecture · Strengths · Weaknesses · Next Steps",
             Inches(0.4), Inches(4.25), Inches(12), Inches(0.5),
             font_size=16, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

    # dot row decoration
    for i in range(20):
        d = slide.shapes.add_shape(9,  # oval
            Inches(0.5 + i * 0.62), Inches(5.6), Inches(0.18), Inches(0.18))
        d.fill.solid()
        col = [ACCENT, ACCENT2, ACCENT3, ACCENT4][i % 4]
        d.fill.fore_color.rgb = col
        d.line.fill.background()


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "What Is RAG-9000?",
                subtitle="A modular, config-driven pipeline for benchmarking retrieval strategies")

    cards = [
        (ACCENT,  "Goal",
         "Evaluate how different retrieval pipelines perform on custom corpora — dense, sparse, and hybrid."),
        (ACCENT2, "Core Idea",
         "A PipelineConfig fully describes one strategy. Six presets cover the main combinations out of the box."),
        (ACCENT3, "Output",
         "Side-by-side metric table (Recall, Precision, MRR, NDCG, MAP, Hit Rate) plus CSV / JSON exports."),
        (ACCENT4, "Use Cases",
         "RAG system selection, retrieval ablations, dataset benchmarking, academic IR experiments."),
    ]

    for i, (col, label, desc) in enumerate(cards):
        row, c = divmod(i, 2)
        x = Inches(0.45 + c * 6.45)
        y = Inches(1.55 + row * 2.55)
        w, h = Inches(6.0), Inches(2.35)

        add_rect(slide, x, y, w, h, fill=BOX_BG, line_color=col)
        # colour tab
        tab = slide.shapes.add_shape(1, x, y, w, Inches(0.38))
        tab.fill.solid()
        tab.fill.fore_color.rgb = col
        tab.line.fill.background()
        add_text(slide, label, x + Inches(0.12), y + Inches(0.03),
                 w - Inches(0.2), Inches(0.34),
                 font_size=13, bold=True, color=DARK_BG)
        add_text(slide, desc, x + Inches(0.18), y + Inches(0.5),
                 w - Inches(0.3), Inches(1.7),
                 font_size=12.5, color=LIGHT_GREY, wrap=True)


def slide_folder_structure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Folder Structure & File Responsibilities")

    files = [
        ("config.py",                   "PipelineConfig dataclass + 6 preset configurations"),
        ("run_evaluation.py",           "CLI entry-point — loads data, builds pipelines, runs Evaluator"),
        ("data/document_store.py",      "In-memory corpus; O(1) ID↔index lookups"),
        ("data/csv_to_json.py",         "CSV → documents.json converter (pandas)"),
        ("data/load_hf_dataset.py",     "HuggingFace dataset → documents + ground truth"),
        ("retrieval/base.py",           "Abstract BaseRetriever / BaseReranker + RetrievalResult"),
        ("retrieval/bi_encoder.py",     "Dense retrieval — FAISS index + sentence-transformers"),
        ("retrieval/bm25_retriever.py", "Sparse retrieval — BM25Okapi"),
        ("retrieval/rrf.py",            "Reciprocal Rank Fusion (merge two result lists)"),
        ("retrieval/cross_encoder.py",  "Cross-encoder reranker (second-stage scoring)"),
        ("retrieval/pipeline.py",       "RetrieverPipeline orchestrator + build_pipeline() factory"),
        ("evaluation/metrics.py",       "Recall, Precision, MRR, NDCG, MAP, HitRate computation"),
        ("evaluation/evaluator.py",     "Evaluator — runs all pipelines, aggregates, prints table"),
    ]

    col_w   = Inches(3.8)
    val_w   = Inches(8.5)
    row_h   = Inches(0.38)
    left_x  = Inches(0.4)
    val_x   = Inches(4.3)
    start_y = Inches(1.45)

    colors = [ACCENT, ACCENT2, ACCENT3, ACCENT4]

    for i, (fname, desc) in enumerate(files):
        y = start_y + Inches(i * 0.395)
        c = colors[i % len(colors)]
        add_rect(slide, left_x, y, col_w, row_h - Inches(0.03), fill=BOX_BG, line_color=c)
        add_text(slide, fname, left_x + Inches(0.1), y + Inches(0.04),
                 col_w - Inches(0.15), row_h,
                 font_size=11, bold=True, color=c)
        add_text(slide, desc, val_x, y + Inches(0.05),
                 val_w, row_h,
                 font_size=11, color=LIGHT_GREY)


def slide_data_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "End-to-End Data Flow")

    steps = [
        (ACCENT4,  "Input",        "CSV file  ·  HuggingFace dataset  ·  Custom JSON"),
        (ACCENT,   "DocumentStore","In-memory corpus — parallel lists + O(1) id↔index hash map"),
        (ACCENT2,  "Stage 1 – Retrieve", "BiEncoderRetriever (FAISS)  ⊕  BM25Retriever (Okapi)"),
        (ACCENT3,  "Stage 2 – Fuse",     "Reciprocal Rank Fusion merges ranked lists (optional)"),
        (ACCENT,   "Stage 3 – Rerank",   "CrossEncoderReranker re-scores (query, doc) pairs (optional)"),
        (ACCENT2,  "Evaluator",    "Computes Recall / Precision / MRR / NDCG / MAP / HitRate per query"),
        (ACCENT4,  "Output",       "Console table  ·  CSV (per-query)  ·  JSON (aggregated)"),
    ]

    box_w  = Inches(11.5)
    box_h  = Inches(0.6)
    left_x = Inches(0.9)
    arr_x  = Inches(6.5)

    for i, (col, label, desc) in enumerate(steps):
        y = Inches(1.35 + i * 0.73)
        add_rect(slide, left_x, y, box_w, box_h, fill=BOX_BG, line_color=col)
        # left label chip
        chip = slide.shapes.add_shape(1, left_x, y, Inches(2.2), box_h)
        chip.fill.solid()
        chip.fill.fore_color.rgb = col
        chip.line.fill.background()
        add_text(slide, label, left_x + Inches(0.1), y + Inches(0.12),
                 Inches(2.0), Inches(0.38),
                 font_size=12, bold=True, color=DARK_BG)
        add_text(slide, desc, left_x + Inches(2.35), y + Inches(0.13),
                 box_w - Inches(2.5), Inches(0.38),
                 font_size=12, color=LIGHT_GREY)
        # arrow between boxes
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(1,
                arr_x, y + box_h,
                Inches(0.03), Inches(0.13))
            arr.fill.solid()
            arr.fill.fore_color.rgb = MID_GREY
            arr.line.fill.background()


def slide_pipeline_presets(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Six Preset Pipeline Configurations",
                subtitle="All defined in config.py — mix and match retrievers, fusion, and reranking")

    headers = ["Pipeline", "Retriever(s)", "Fusion", "Reranker", "Best For"]
    rows = [
        ("bi_encoder_only",      "FAISS (dense)",     "—",   "—",            "Fast semantic search"),
        ("bm25_only",            "BM25 (sparse)",     "—",   "—",            "Keyword / no-GPU setup"),
        ("bi_encoder_reranked",  "FAISS",             "—",   "Cross-Encoder","Dense + accuracy boost"),
        ("bm25_reranked",        "BM25",              "—",   "Cross-Encoder","Keyword + accuracy boost"),
        ("hybrid_rrf",           "FAISS + BM25",      "RRF", "—",            "Balanced, no reranker cost"),
        ("hybrid_rrf_reranked",  "FAISS + BM25",      "RRF", "Cross-Encoder","Maximum accuracy"),
    ]

    col_xs  = [Inches(0.35), Inches(3.2), Inches(5.65), Inches(7.1), Inches(8.9)]
    col_ws  = [Inches(2.75), Inches(2.3), Inches(1.35), Inches(1.7), Inches(4.1)]
    header_y = Inches(1.5)
    row_h    = Inches(0.55)
    hdr_cols = [ACCENT, ACCENT2, ACCENT3, ACCENT4, ACCENT]

    # header row
    for j, (hdr, hx, hw, hc) in enumerate(zip(headers, col_xs, col_ws, hdr_cols)):
        add_rect(slide, hx, header_y, hw, row_h - Inches(0.04), fill=hc)
        add_text(slide, hdr, hx + Inches(0.08), header_y + Inches(0.1),
                 hw - Inches(0.1), row_h,
                 font_size=12, bold=True, color=DARK_BG)

    row_fills = [BOX_BG, RGBColor(0x28, 0x28, 0x40)]
    for i, row in enumerate(rows):
        y = header_y + row_h + Inches(i * 0.55)
        for j, (val, hx, hw) in enumerate(zip(row, col_xs, col_ws)):
            fill = row_fills[i % 2]
            add_rect(slide, hx, y, hw, row_h - Inches(0.04), fill=fill)
            txt_col = ACCENT if j == 0 else (ACCENT2 if val not in ("—", "") else MID_GREY)
            add_text(slide, val, hx + Inches(0.08), y + Inches(0.1),
                     hw - Inches(0.1), row_h,
                     font_size=11.5, color=txt_col, bold=(j == 0))

    # footnote
    add_text(slide, "top_k_retrieve=50  ·  top_k_final=10  ·  rrf_k=60  (all overrideable via PipelineConfig)",
             Inches(0.35), Inches(7.1), Inches(12.5), Inches(0.35),
             font_size=10, color=MID_GREY, italic=True)


def slide_retrieval_components(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Retrieval Components Deep-Dive")

    cards = [
        (ACCENT, "BiEncoderRetriever",
         ["sentence-transformers model (all-MiniLM-L6-v2 default)",
          "FAISS index — Inner Product or L2 distance",
          "Vectorised batch_retrieve() for efficiency",
          "L2-norm flag + score sign fix for IP vs L2"]),
        (ACCENT2, "BM25Retriever",
         ["BM25Okapi (rank-bm25) — TF-IDF + length norm",
          "Tokenise at init time (lowercase + split on \\W)",
          "Stops at score 0.0 to skip zero-match docs",
          "No GPU dependency — runs anywhere"]),
        (ACCENT3, "RRF Fusion",
         ["Score = Σ 1 / (k + rank_i(doc)) across lists",
          "Rank-based: immune to score magnitude mismatches",
          "Single hyperparameter k=60 (standard default)",
          "Documents in multiple lists get natural boost"]),
        (ACCENT4, "CrossEncoderReranker",
         ["ms-marco-MiniLM-L-6-v2 default",
          "Scores (query, passage) pairs jointly",
          "Applied only to small candidate pool (~50)",
          "Batch inference (batch_size=32) for speed"]),
    ]

    for i, (col, title, bullets) in enumerate(cards):
        row, c = divmod(i, 2)
        x = Inches(0.4 + c * 6.45)
        y = Inches(1.55 + row * 2.65)
        w, h = Inches(6.0), Inches(2.5)

        add_rect(slide, x, y, w, h, fill=BOX_BG, line_color=col)
        tab = slide.shapes.add_shape(1, x, y, w, Inches(0.38))
        tab.fill.solid()
        tab.fill.fore_color.rgb = col
        tab.line.fill.background()
        add_text(slide, title, x + Inches(0.12), y + Inches(0.03),
                 w - Inches(0.2), Inches(0.35),
                 font_size=13, bold=True, color=DARK_BG)
        for j, b in enumerate(bullets):
            add_text(slide, "▸  " + b,
                     x + Inches(0.18), y + Inches(0.5 + j * 0.44),
                     w - Inches(0.3), Inches(0.42),
                     font_size=11.5, color=LIGHT_GREY)


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Evaluation Metrics",
                subtitle="All computed at rank K (default 10) — binary relevance")

    metrics = [
        ("Recall@K",    ACCENT,  "hits / n_relevant",
         "Fraction of all relevant docs found in top-K.  High = broad coverage."),
        ("Precision@K", ACCENT2, "hits / K",
         "Fraction of the top-K results that are relevant.  High = low false-positive rate."),
        ("MRR",         ACCENT3, "1 / rank(first relevant)",
         "Rewards finding the first relevant result as early as possible."),
        ("NDCG@K",      ACCENT4, "DCG / Ideal_DCG",
         "Position-weighted score.  Penalises relevant docs ranked too low.  Primary metric."),
        ("MAP@K",       ACCENT,  "Σ P@hit / min(nRel, K)",
         "Average precision at each relevant position.  Balanced recall-precision summary."),
        ("Hit Rate@K",  ACCENT2, "1 if any hit else 0",
         "Binary success metric — did we find at least one relevant doc in top-K?"),
    ]

    for i, (name, col, formula, interp) in enumerate(metrics):
        row, c = divmod(i, 2)
        x = Inches(0.4 + c * 6.45)
        y = Inches(1.55 + row * 1.82)
        w, h = Inches(6.0), Inches(1.65)

        add_rect(slide, x, y, w, h, fill=BOX_BG, line_color=col)
        add_text(slide, name, x + Inches(0.15), y + Inches(0.1),
                 Inches(2.5), Inches(0.38),
                 font_size=14, bold=True, color=col)
        add_text(slide, formula, x + Inches(2.7), y + Inches(0.12),
                 Inches(3.1), Inches(0.35),
                 font_size=11.5, color=ACCENT4, italic=True)
        add_text(slide, interp, x + Inches(0.15), y + Inches(0.55),
                 w - Inches(0.3), Inches(1.0),
                 font_size=11.5, color=LIGHT_GREY, wrap=True)


def slide_strengths(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Strengths", subtitle="What the framework does well")

    strengths = [
        "Fully modular — retrievers, rerankers, fusion are all independently swappable",
        "Config-driven design: one PipelineConfig object fully describes a strategy",
        "Six ready-made presets cover the most common retrieval combinations",
        "Comprehensive metrics: Recall, Precision, MRR, NDCG, MAP, Hit Rate — all in one run",
        "Efficient batch encoding for bi-encoder (vectorised sentence-transformers)",
        "RRF fusion requires no score normalisation — robust to dense/sparse magnitude gaps",
        "Multiple output formats: live table, per-query CSV, aggregated JSON",
        "Multiple data sources: CSV, HuggingFace datasets, custom JSON",
        "Clear three-layer separation: DocumentStore → Pipeline → Evaluator",
        "Strategy pattern with abstract base classes — easy to extend",
        "Factory function (build_pipeline) decouples config from instantiation",
        "O(1) ID ↔ index lookups in DocumentStore via hash map",
    ]

    n = len(strengths)
    half = (n + 1) // 2
    left_items  = strengths[:half]
    right_items = strengths[half:]

    col_w  = Inches(6.0)
    gap    = Inches(0.5)
    left_x = Inches(0.45)
    right_x = left_x + col_w + gap
    start_y = Inches(1.5)
    row_h  = Inches(0.43)

    for i, item in enumerate(left_items):
        y = start_y + Inches(i * 0.455)
        add_rect(slide, left_x, y, col_w, row_h, fill=BOX_BG, line_color=ACCENT2)
        add_text(slide, "✔  " + item, left_x + Inches(0.1), y + Inches(0.05),
                 col_w - Inches(0.15), row_h,
                 font_size=11.5, color=WHITE)

    for i, item in enumerate(right_items):
        y = start_y + Inches(i * 0.455)
        add_rect(slide, right_x, y, col_w, row_h, fill=BOX_BG, line_color=ACCENT2)
        add_text(slide, "✔  " + item, right_x + Inches(0.1), y + Inches(0.05),
                 col_w - Inches(0.15), row_h,
                 font_size=11.5, color=WHITE)


def slide_weaknesses(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Weaknesses & Limitations", subtitle="Known gaps and design trade-offs")

    weaknesses = [
        ("FAISS index built externally",       "No index-building code; alignment with DocumentStore is manual and error-prone"),
        ("Binary relevance only",              "No graded relevance support — NDCG gains are always 1.0 for any hit"),
        ("Single K value per run",             "Evaluating at multiple cutoffs (5, 10, 20) requires separate runs"),
        ("Simple BM25 tokenisation",           "Lowercase + whitespace split only — no stopwords, stemming, or lemmatisation"),
        ("No query preprocessing",             "No spell-correction, expansion, or normalisation before retrieval"),
        ("No metadata / field-based filters",  "Can't restrict retrieval by date, category, or other metadata fields"),
        ("No distributed evaluation",          "All queries processed sequentially — large query sets will be slow"),
        ("No component-level ablation",        "Hard to attribute poor NDCG to retrieval vs. reranking stage"),
        ("Cross-encoder reranker only",        "No listwise LTR, ColBERT late-interaction, or learned fusion options"),
        ("Model management is implicit",       "HuggingFace models auto-downloaded — can fail offline or in CI/CD"),
    ]

    col_w   = Inches(2.8)
    desc_w  = Inches(9.4)
    left_x  = Inches(0.4)
    desc_x  = Inches(3.35)
    start_y = Inches(1.5)
    row_h   = Inches(0.46)
    fills   = [BOX_BG, RGBColor(0x28, 0x28, 0x40)]

    for i, (label, desc) in enumerate(weaknesses):
        y = start_y + Inches(i * 0.49)
        f = fills[i % 2]
        add_rect(slide, left_x, y, col_w, row_h, fill=f, line_color=ACCENT3)
        add_text(slide, label, left_x + Inches(0.1), y + Inches(0.06),
                 col_w - Inches(0.15), row_h,
                 font_size=11.5, bold=True, color=ACCENT3)
        add_rect(slide, desc_x, y, desc_w, row_h, fill=f)
        add_text(slide, desc, desc_x + Inches(0.1), y + Inches(0.06),
                 desc_w - Inches(0.15), row_h,
                 font_size=11.5, color=LIGHT_GREY)


def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    slide_title(slide, "Suggested Next Steps",
                subtitle="Prioritised improvements to expand capability and robustness")

    steps = [
        (ACCENT,  "P1 – Quick Wins",
         [
             "Add index-building script (sentence-transformers → FAISS) to remove external dependency",
             "Multi-K evaluation in a single run (e.g. @5, @10, @20) for richer analysis",
             "Add minimum-score threshold filter to suppress near-zero BM25 results",
         ]),
        (ACCENT2, "P2 – Retrieval Quality",
         [
             "Upgrade BM25 tokeniser: stopword removal + stemming (NLTK / spaCy)",
             "Add query preprocessing: lowercasing, spell-correction, synonym expansion",
             "Support graded relevance levels (0/1/2) for proper NDCG computation",
         ]),
        (ACCENT3, "P3 – Architecture Extensibility",
         [
             "Add ColBERT late-interaction retriever as a third retriever option",
             "Add a learned fusion option (e.g. linear interpolation with tunable α)",
             "Component-level metrics: log per-stage recall to isolate retrieval vs. reranking gaps",
         ]),
        (ACCENT4, "P4 – Production Readiness",
         [
             "Parallelise query evaluation (multiprocessing / asyncio) for large query sets",
             "Explicit model registry / pin model versions for offline / CI reproducibility",
             "Metadata-aware pre-filtering (date range, category) before retrieval stage",
         ]),
    ]

    for i, (col, label, bullets) in enumerate(steps):
        row, c = divmod(i, 2)
        x = Inches(0.4 + c * 6.45)
        y = Inches(1.55 + row * 2.7)
        w, h = Inches(6.0), Inches(2.55)

        add_rect(slide, x, y, w, h, fill=BOX_BG, line_color=col)
        tab = slide.shapes.add_shape(1, x, y, w, Inches(0.38))
        tab.fill.solid()
        tab.fill.fore_color.rgb = col
        tab.line.fill.background()
        add_text(slide, label, x + Inches(0.12), y + Inches(0.03),
                 w - Inches(0.2), Inches(0.35),
                 font_size=12, bold=True, color=DARK_BG)
        for j, b in enumerate(bullets):
            add_text(slide, "▸  " + b,
                     x + Inches(0.18), y + Inches(0.5 + j * 0.65),
                     w - Inches(0.3), Inches(0.62),
                     font_size=11.5, color=LIGHT_GREY, wrap=True)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)

    # large decorative circle
    circ = slide.shapes.add_shape(9, Inches(4.16), Inches(1.5), Inches(5.0), Inches(5.0))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BOX_BG
    circ.line.color.rgb = ACCENT
    circ.line.width = Pt(2)

    add_text(slide, "RAG-9000", Inches(3), Inches(2.7), Inches(7.33), Inches(1.0),
             font_size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "Modular  ·  Extensible  ·  Comprehensive",
             Inches(3), Inches(3.65), Inches(7.33), Inches(0.5),
             font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "A solid foundation for systematic retrieval evaluation",
             Inches(2.0), Inches(4.2), Inches(9.33), Inches(0.45),
             font_size=13, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

    add_accent_bar(slide, color=ACCENT2, height=Inches(0.07))
    # bottom bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT2
    bar.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_folder_structure(prs)
    slide_data_flow(prs)
    slide_pipeline_presets(prs)
    slide_retrieval_components(prs)
    slide_metrics(prs)
    slide_strengths(prs)
    slide_weaknesses(prs)
    slide_next_steps(prs)
    slide_closing(prs)

    out = r"c:\Users\mathis.de.greef\OneDrive - Accenture\Desktop\rag-9000\RAG-9000_Architecture.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
